"""
Generate a weekly report PPTX and/or PDF from report form data.

Usage:
    python scripts/generate_pptx_report.py --data <json_file> --output <pptx_file> [--pdf <pdf_file>]

The JSON file contains:
    template_path: str        — path to the original PPTX template (only needed for PPTX)
    report: object            — WeeklyReport document
      weekStart, weekEnd, sections[], summary

Charts and tables use data from report.sections[] (form input) only.
Each section has values with headingKey, previousValue, currentValue, targetValue.
"""

import io
import json
import os
import re
import sys
import argparse
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.patches as mpatches
import numpy as np

import plotly.graph_objects as go
from plotly.subplots import make_subplots

from fpdf import FPDF

FONT_DIR = Path(__file__).resolve().parent / 'fonts'

COLORS = {
    'primary':       '#B41531',
    'primary_light': '#EF4444',
    'secondary':     '#DC2626',
    'accent':        '#7F1D1D',
    'bg_card':       '#F8FAFC',
    'card_border':   '#E2E8F0',
    'bar_fill':      '#B41531',
    'bar_prev':      '#94A3B8',
    'bar_target':    '#EA580C',
    'text_dark':     '#1F1F24',
    'text_muted':    '#64748B',
    'grid':          '#E2E8F0',
    'grid_major':    '#CBD5E1',
    'pct_pos_bg':    '#DCEFDD',
    'pct_pos_text':  '#1A7F37',
    'pct_neg_bg':    '#FDE8E8',
    'pct_neg_text':  '#C0392B',
}

FONT_SIZES = {
    'title': 18,
    'bar_label': 12,
    'axis_label': 9,
    'tick': 10,
    'target_note': 7,
    'pct_change': 13,
}


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _transparent_png():
    """Return a 1×1 transparent PNG BytesIO for zero-value charts."""
    import struct, zlib
    buf = io.BytesIO()
    raw = b'\x00' * 4
    def chunk(ctype, data):
        c = ctype + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    buf.write(b'\x89PNG\r\n\x1a\n')
    buf.write(chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 6, 0, 0, 0)))
    buf.write(chunk(b'IDAT', zlib.compress(raw)))
    buf.write(chunk(b'IEND', b''))
    buf.seek(0)
    return buf


def _lighten_color(hex_color, factor=1.25):
    """Lighten a hex color. factor > 1 = lighter, < 1 = darker."""
    from matplotlib.colors import to_rgb, to_hex
    r, g, b = to_rgb(hex_color)
    r = min(1.0, r * factor)
    g = min(1.0, g * factor)
    b = min(1.0, b * factor)
    return to_hex((r, g, b))


def _text_callout(message, sub_message='', fig_width=7.0, fig_height=4.3):
    """Render a centered text callout image for zero-value or special cases."""
    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    fig.patch.set_alpha(0)
    ax.set_facecolor('#FFFFFF')
    ax.axis('off')
    ax.text(0.5, 0.55, message, ha='center', va='center',
            fontsize=24, fontweight='bold', color='#B41531',
            transform=ax.transAxes, fontfamily='sans-serif')
    if sub_message:
        ax.text(0.5, 0.38, sub_message, ha='center', va='center',
                fontsize=14, color='#64748B',
                transform=ax.transAxes, fontfamily='sans-serif')
    plt.tight_layout(pad=0.8)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=300, facecolor='none', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


def create_bar_chart(values, labels, title, target=None, dpi=300, chart_type='bar', prev_pct=None, fig_width=None, fig_height=None,
                     y_min=None, y_max=None):
    """Create a premium, modern bar chart with brand-red and clean silver-gray palette."""

    if target is not None and len(values) == 2:
        values = list(values) + [target]
        labels = list(labels) + ['Target']

    if all(v == 0 for v in values):
        return _text_callout('0 Products Rejected', '100% Pass Rate', fig_width or 7.0, fig_height or 4.3)

    if prev_pct is None and len(values) >= 2 and values[0] > 0:
        prev_pct = round(((values[1] - values[0]) / values[0]) * 100, 1)

    is_horiz = chart_type == 'horizontal'
    if fig_width is None or fig_height is None:
        fig_width, fig_height = (7.2, 4.2) if not is_horiz else (6.4, 3.6)

    # Convert inches to pixels
    width_px = int(fig_width * dpi)
    height_px = int(fig_height * dpi)

    # Determine bar colors
    bar_colors = []
    for i in range(len(values)):
        if i == 0:
            bar_colors.append('#CBD5E1')  # Previous week
        elif i == 1:
            bar_colors.append('#B41531')  # Current week
        else:
            bar_colors.append('white')    # Target (will have red border)

    fig = go.Figure()

    # Add bars
    fig.add_trace(go.Bar(
        x=list(range(len(values))) if not is_horiz else values,
        y=values if not is_horiz else list(range(len(values))),
        orientation='h' if is_horiz else 'v',
        marker_color=bar_colors,
        marker_line_color='#B41531' if not is_horiz else '#B41531',
        marker_line_width=1.8,
        width=0.45,
        text=[f'{v:,}' for v in values],
        textposition='outside',
        hoverinfo='skip'
    ))

    # Update layout
    fig.update_layout(
        title_text=title,
        title_x=0.5,
        title_font=dict(size=FONT_SIZES['title'], family='sans-serif'),
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='sans-serif', size=FONT_SIZES['tick'], color=COLORS['text_muted']),
        showlegend=False,
        margin=dict(l=20, r=20, t=40, b=20)
    )

    # Configure axes
    if is_horiz:
        fig.update_yaxes(
            tickmode='array',
            tickvals=list(range(len(values))),
            ticktext=labels,
            tickfont=dict(size=FONT_SIZES['tick'], family='sans-serif', color=COLORS['text_muted']),
            autorange="reversed",  # to match original order (top to bottom)
            showgrid=False,
            zeroline=False
        )
        fig.update_xaxes(
            showgrid=True,
            gridcolor='#F1F5F9',
            zeroline=False,
            tickformat=',',
            tickfont=dict(size=FONT_SIZES['tick'], family='sans-serif', color=COLORS['text_muted']),
            rangemode='tozero'
        )
        if y_min is not None and y_max is not None:
            fig.update_xaxes(range=[y_min, y_max])
        else:
            fig.update_xaxes(range=[0, max(values) * 1.25])
    else:
        fig.update_xaxes(
            tickmode='array',
            tickvals=list(range(len(values))),
            ticktext=labels,
            tickfont=dict(size=FONT_SIZES['tick'], family='sans-serif', color=COLORS['text_muted']),
            showgrid=False,
            zeroline=False
        )
        fig.update_yaxes(
            showgrid=True,
            gridcolor='#F1F5F9',
            zeroline=False,
            tickformat=',',
            tickfont=dict(size=FONT_SIZES['tick'], family='sans-serif', color=COLORS['text_muted']),
            rangemode='tozero'
        )
        if y_min is not None and y_max is not None:
            fig.update_yaxes(range=[y_min, y_max])
        else:
            fig.update_yaxes(range=[0, max(values) * 1.25])

    # Add percentage change badge as annotation
    if prev_pct is not None and len(values) >= 2:
        pct_sym = '+' if prev_pct >= 0 else ''
        pct_text = f'{pct_sym}{prev_pct:.1f}%'
        # Position: top center of the chart
        fig.add_annotation(
            xref="paper", yref="paper",
            x=0.5, y=1.12,
            text=pct_text,
            showarrow=False,
            font=dict(size=FONT_SIZES['pct_change'], family='sans-serif', color='white'),
            bgcolor='#C0392B',
            borderpad=4
        )

    # Generate PNG image bytes
    img_bytes = fig.to_image(format="png", width=width_px, height=height_px, scale=1)
    buf = io.BytesIO(img_bytes)
    return buf


def create_trend_chart(values, labels, title, dpi=300, fig_width=None, fig_height=None, y_min=None, y_max=None):
    """Create a premium, modern trend line/area chart with brand-red theme."""
    if fig_width is None or fig_height is None:
        fig_width, fig_height = 7.2, 4.2

    # Convert inches to pixels
    width_px = int(fig_width * dpi)
    height_px = int(fig_height * dpi)

    # Clean the values and ensure numeric
    y = [v if v is not None else 0 for v in values]
    max_val = max(y) if len(y) > 0 else 1
    if max_val == 0:
        max_val = 1

    n = len(y)
    x = list(range(n))

    fig = go.Figure()

    # Add area and line
    fig.add_trace(go.Scatter(
        x=x + x[::-1],  # x, then x reversed for filled area
        y=y + [0]*len(y),
        fill='toself',
        fillcolor='rgba(180, 21, 49, 0.08)',  # brand red with alpha 0.08
        line=dict(color='rgba(180,21,49,0)'),  # invisible line for area
        hoverinfo='skip',
        showlegend=False
    ))
    fig.add_trace(go.Scatter(
        x=x,
        y=y,
        mode='lines+markers',
        line=dict(color='#B41531', width=3),
        marker=dict(
            size=[5 if i < n-1 else 8 for i in range(n)],  # smaller for intermediate, larger for last
            color=['white' if i < n-1 else '#B41531' for i in range(n)],
            line=dict(color='#B41531', width=2)
        ),
        hoverinfo='skip',
        showlegend=False
    ))

    # Update layout
    fig.update_layout(
        title_text=title,
        title_x=0.5,
        title_font=dict(size=FONT_SIZES['title'], family='sans-serif'),
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='sans-serif', size=FONT_SIZES['tick'], color=COLORS['text_muted']),
        showlegend=False,
        margin=dict(l=20, r=20, t=40, b=20)
    )

    # Configure axes
    fig.update_xaxes(
        tickmode='array',
        tickvals=x,
        ticktext=labels,
        tickfont=dict(size=10, family='sans-serif', color=COLORS['text_muted']),
        showgrid=False,
        zeroline=False
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor='#F1F5F9',
        zeroline=False,
        tickformat=',',
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        rangemode='tozero'
    )

    # Set limits
    fig.update_xaxes(range=[-0.2, n - 0.8 if n > 1 else 0.5])
    if y_min is not None and y_max is not None:
        fig.update_yaxes(range=[y_min, y_max])
    else:
        fig.update_yaxes(range=[0, max_val * 1.2])

    # Add data labels on points
    for xi, yi in zip(x, y):
        fig.add_annotation(
            x=xi, y=yi + max_val * 0.04,
            text=f'{int(yi):,}',
            showarrow=False,
            font=dict(size=10, family='sans-serif', color='#1F1F24'),
            yanchor='bottom'
        )

    # Generate PNG image bytes
    img_bytes = fig.to_image(format="png", width=width_px, height=height_px, scale=1)
    buf = io.BytesIO(img_bytes)
    return buf


def create_kpi_trend_chart(series_data, labels, title, dpi=300, fig_width=None, fig_height=None):
    """Create a premium, multi-series trend line chart for KPI Overview."""
    if fig_width is None or fig_height is None:
        fig_width, fig_height = 6.4, 3.6

    # Convert inches to pixels
    width_px = int(fig_width * dpi)
    height_px = int(fig_height * dpi)

    # Cohesive red/white compatible color palette
    colors = ['#B41531', '#EF4444', '#475569']  # Brand Red, Coral, Contrast Slate
    
    max_val = 1
    series_dict = {}
    for idx, (name, values) in enumerate(series_data.items()):
        y = [v if v is not None else 0 for v in values]
        max_val = max(max_val, max(y))
        series_dict[name] = y
        col = colors[idx % len(colors)]
        series_dict[name + '_color'] = col

    fig = go.Figure()

    for idx, (name, values) in enumerate(series_data.items()):
        y = series_dict[name]
        col = series_dict[name + '_color']
        x = list(range(len(y)))
        
        # Add area fill
        fig.add_trace(go.Scatter(
            x=x + x[::-1],
            y=y + [0]*len(y),
            fill='toself',
            fillcolor=f'rgba{tuple(list(hex_to_rgb(col)) + [0.03])}',
            line=dict(color='rgba(255,255,255,0)'),
            hoverinfo='skip',
            showlegend=False
        ))
        # Add line and markers
        fig.add_trace(go.Scatter(
            x=x,
            y=y,
            mode='lines+markers',
            name=name,
            line=dict(color=col, width=2.5),
            marker=dict(
                size=[8 if i == len(y)-1 else 6 for i in range(len(y))],
                color='white',
                line=dict(color=col, width=1.8)
            ),
            hoverinfo='skip',
            showlegend=True
        ))
        # Add data label on last point
        if len(y) > 0:
            fig.add_annotation(
                x=len(y)-1,
                y=y[-1] + max_val * 0.03,
                text=f'{int(y[-1]):,}',
                showarrow=False,
                font=dict(size=9, family='sans-serif', color=col),
                yanchor='bottom'
            )

    # Update layout
    fig.update_layout(
        title_text=title,
        title_x=0.5,
        title_font=dict(size=FONT_SIZES['title'], family='sans-serif'),
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='sans-serif', size=FONT_SIZES['tick'], color=COLORS['text_muted']),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1
        ),
        margin=dict(l=20, r=20, t=40, b=20)
    )

    # Configure axes
    fig.update_xaxes(
        tickmode='array',
        tickvals=list(range(len(labels))),
        ticktext=labels,
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        showgrid=False,
        zeroline=False
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor='#F1F5F9',
        zeroline=False,
        tickformat=',',
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        rangemode='tozero'
    )

    # Set limits
    fig.update_xaxes(range=[-0.2, len(labels) - 0.5 if len(labels) > 1 else 0.5])
    fig.update_yaxes(range=[0, max_val * 1.25])

    # Generate PNG image bytes
    img_bytes = fig.to_image(format="png", width=width_px, height=height_px, scale=1)
    buf = io.BytesIO(img_bytes)
    return buf

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_trend_chart(values, labels, title, dpi=300, fig_width=None, fig_height=None, y_min=None, y_max=None):
    """Create a premium, modern trend line/area chart with brand-red theme."""
    if fig_width is None or fig_height is None:
        fig_width, fig_height = 7.2, 4.2

    # Convert inches to pixels
    width_px = int(fig_width * dpi)
    height_px = int(fig_height * dpi)

    # Clean the values and ensure numeric
    y = [v if v is not None else 0 for v in values]
    max_val = max(y) if len(y) > 0 else 1
    if max_val == 0:
        max_val = 1

    n = len(y)
    x = list(range(n))

    fig = go.Figure()

    # Add area and line
    fig.add_trace(go.Scatter(
        x=x + x[::-1],  # x, then x reversed for filled area
        y=y + [0]*len(y),
        fill='toself',
        fillcolor='rgba(180, 21, 49, 0.08)',  # brand red with alpha 0.08
        line=dict(color='rgba(180,21,49,0)'),  # invisible line for area
        hoverinfo='skip',
        showlegend=False
    ))
    fig.add_trace(go.Scatter(
        x=x,
        y=y,
        mode='lines+markers',
        line=dict(color='#B41531', width=3),
        marker=dict(
            size=[5 if i < n-1 else 8 for i in range(n)],  # smaller for intermediate, larger for last
            color=['white' if i < n-1 else '#B41531' for i in range(n)],
            line=dict(color='#B41531', width=2)
        ),
        hoverinfo='skip',
        showlegend=False
    ))

    # Update layout
    fig.update_layout(
        title_text=title,
        title_x=0.5,
        title_font=dict(size=FONT_SIZES['title'], family='sans-serif'),
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='sans-serif', size=FONT_SIZES['tick'], color=COLORS['text_muted']),
        showlegend=False,
        margin=dict(l=20, r=20, t=40, b=20)
    )

    # Configure axes
    fig.update_xaxes(
        tickmode='array',
        tickvals=x,
        ticktext=labels,
        tickfont=dict(size=10, family='sans-serif', color=COLORS['text_muted']),
        showgrid=False,
        zeroline=False
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor='#F1F5F9',
        zeroline=False,
        tickformat=',',
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        rangemode='tozero'
    )

    # Set limits
    fig.update_xaxes(range=[-0.2, n - 0.8 if n > 1 else 0.5])
    if y_min is not None and y_max is not None:
        fig.update_yaxes(range=[y_min, y_max])
    else:
        fig.update_yaxes(range=[0, max_val * 1.2])

    # Add data labels on points
    for xi, yi in zip(x, y):
        fig.add_annotation(
            x=xi, y=yi + max_val * 0.04,
            text=f'{int(yi):,}',
            showarrow=False,
            font=dict(size=10, family='sans-serif', color='#1F1F24'),
            yanchor='bottom'
        )

    # Generate PNG image bytes
    img_bytes = fig.to_image(format="png", width=width_px, height=height_px, scale=1)
    buf = io.BytesIO(img_bytes)
    return buf


def create_kpi_trend_chart(series_data, labels, title, dpi=300, fig_width=None, fig_height=None):
    """Create a premium, multi-series trend line chart for KPI Overview."""
    if fig_width is None or fig_height is None:
        fig_width, fig_height = 6.4, 3.6

    # Convert inches to pixels
    width_px = int(fig_width * dpi)
    height_px = int(fig_height * dpi)

    # Cohesive red/white compatible color palette
    colors = ['#B41531', '#EF4444', '#475569']  # Brand Red, Coral, Contrast Slate
    
    max_val = 1
    series_dict = {}
    for idx, (name, values) in enumerate(series_data.items()):
        y = [v if v is not None else 0 for v in values]
        max_val = max(max_val, max(y))
        series_dict[name] = y
        col = colors[idx % len(colors)]
        series_dict[name + '_color'] = col

    fig = go.Figure()

    for idx, (name, values) in enumerate(series_data.items()):
        y = series_dict[name]
        col = series_dict[name + '_color']
        x = list(range(len(y)))
        
        # Add area fill
        fig.add_trace(go.Scatter(
            x=x + x[::-1],
            y=y + [0]*len(y),
            fill='toself',
            fillcolor=f'rgba{tuple(list(hex_to_rgb(col)) + [0.03])}',
            line=dict(color='rgba(255,255,255,0)'),
            hoverinfo='skip',
            showlegend=False
        ))
        # Add line and markers
        fig.add_trace(go.Scatter(
            x=x,
            y=y,
            mode='lines+markers',
            name=name,
            line=dict(color=col, width=2.5),
            marker=dict(
                size=[8 if i == len(y)-1 else 6 for i in range(len(y))],
                color='white',
                line=dict(color=col, width=1.8)
            ),
            hoverinfo='skip',
            showlegend=True
        ))
        # Add data label on last point
        if len(y) > 0:
            fig.add_annotation(
                x=len(y)-1,
                y=y[-1] + max_val * 0.03,
                text=f'{int(y[-1]):,}',
                showarrow=False,
                font=dict(size=9, family='sans-serif', color=col),
                yanchor='bottom'
            )

    # Update layout
    fig.update_layout(
        title_text=title,
        title_x=0.5,
        title_font=dict(size=FONT_SIZES['title'], family='sans-serif'),
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='sans-serif', size=FONT_SIZES['tick'], color=COLORS['text_muted']),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1
        ),
        margin=dict(l=20, r=20, t=40, b=20)
    )

    # Configure axes
    fig.update_xaxes(
        tickmode='array',
        tickvals=list(range(len(labels))),
        ticktext=labels,
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        showgrid=False,
        zeroline=False
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor='#F1F5F9',
        zeroline=False,
        tickformat=',',
        tickfont=dict(size=9, family='sans-serif', color=COLORS['text_muted']),
        rangemode='tozero'
    )

    # Set limits
    fig.update_xaxes(range=[-0.2, len(labels) - 0.5 if len(labels) > 1 else 0.5])
    fig.update_yaxes(range=[0, max_val * 1.25])

    # Generate PNG image bytes
    img_bytes = fig.to_image(format="png", width=width_px, height=height_px, scale=1)
    buf = io.BytesIO(img_bytes)
    return buf


def find_picture_by_position(slide, target_left, target_top, tolerance=50000):
    for shape in slide.shapes:
        if hasattr(shape, 'image') and shape.shape_type == 13:
            if abs(shape.left - target_left) < tolerance and abs(shape.top - target_top) < tolerance:
                return shape
    return None


def _replace_image_with_new_part(shape, image_bytes):
    from pptx.oxml.ns import qn
    blip = shape._element.find('.//' + qn('a:blip'))
    if blip is not None:
        rId = blip.get(qn('r:embed'))
        if rId:
            image_part = shape.part.related_part(rId)
            image_part._blob = image_bytes


def _replace_image_with_new_part(shape, image_bytes):
    """Replace image by creating a new image part, avoiding sharing issues with cloned slides."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    from io import BytesIO
    blip = shape._element.find('.//' + qn('a:blip'))
    if blip is None:
        return
    rId_old = blip.get(qn('r:embed'))
    if not rId_old:
        return
    pkg = shape.part.package
    buf = BytesIO(image_bytes)
    new_img_part = pkg.get_or_add_image_part(buf)
    new_rId = shape.part.relate_to(new_img_part, RT.IMAGE)
    blip.set(qn('r:embed'), new_rId)


def format_val(v):
    if v is None:
        return '—'
    if isinstance(v, float):
        if v == int(v):
            return f'{int(v):,}'
        return f'{v:,.1f}'
    return f'{int(v):,}'


LABEL_KEY_MAP = {
    'total vendors': 'totalVendors',
    'verified vendors': 'verifiedVendors',
    'total marketplace products': 'totalMarketplaceProducts',
    'daily average listings': 'dailyAverageListings',
    'backlog': 'backlogProducts',
    'specifications added': 'totalSpecificationsAdded',
    'spec completion': 'specificationCompletionPercent',
    'specification completion': 'specificationCompletionPercent',
    'products approved': 'productsApproved',
    'products rejected': 'productsRejected',
    'products pending': 'productsPending',
}


def find_heading_key(text):
    lower = text.lower().strip()
    for label, key in LABEL_KEY_MAP.items():
        if label in lower:
            return key
    return None


LEFT_COL = {'min': 0, 'max': 2500000}
CENTER_COL = {'min': 2500000, 'max': 6500000}
RIGHT_COL = {'min': 6500000, 'max': 13000000}


def get_column(x):
    if LEFT_COL['min'] <= x <= LEFT_COL['max']:
        return 'prev'
    elif CENTER_COL['min'] <= x <= CENTER_COL['max']:
        return 'current'
    elif RIGHT_COL['min'] <= x <= RIGHT_COL['max']:
        return 'target'
    return None


COLUMN_FIELD_MAP = {
    'prev': 'previousValue',
    'current': 'currentValue',
    'target': 'targetValue',
}

def get_value_from_report(heading_key, column, sections):
    field = COLUMN_FIELD_MAP.get(column)
    if not field:
        return None
    for section in sections:
        for v in section.get('values', []):
            if v.get('headingKey') == heading_key:
                val = v.get(field)
                if val is not None:
                    return val
    return None


def get_section_by_dept(sections, dept_name):
    """Find section by department name (case-insensitive partial match)."""
    for section in sections:
        if dept_name.lower() in section.get('departmentName', '').lower():
            return section
    return None


def get_value_from_section(section, heading_key, column):
    """Get value from a specific section by headingKey and column (prev/current/target/previous)."""
    if not section:
        return None
    field = COLUMN_FIELD_MAP.get(column) or (column + 'Value')
    for v in section.get('values', []):
        if v.get('headingKey') == heading_key:
            return v.get(field)
    return None


def get_chart_data_3bars(sections, dept_name, heading_key):
    """Extract [prev, current, target] for a heading from a department section."""
    section = get_section_by_dept(sections, dept_name)
    if not section:
        return [0, 0, 0]
    prev = get_value_from_section(section, heading_key, 'previous') or 0
    curr = get_value_from_section(section, heading_key, 'current') or 0
    target = get_value_from_section(section, heading_key, 'target') or 0
    return [prev, curr, target]


def replace_text_in_shape(shape, report_sections, summary, date_str):
    if not shape.has_text_frame:
        return False
    modified = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = run.text
            new_text = text

            # Date range replacement
            if '{{date_range}}' in new_text:
                new_text = new_text.replace('{{date_range}}', date_str)

            if new_text != text:
                run.text = new_text
                modified = True
    return modified


def replace_label_values_in_shape(shape, sections, summary, date_str, slide_idx=None, slide_count=9):
    """Replace only heading-key-matched text in column zones; clear unmatched text.

    Date-range replacement runs on ALL shapes (including headers/footers).
    Heading-key replacement only applies to column-zone shapes.
    Stale-text clearing only applies on data-table slides (1, 4, 6, and 10 in new template).
    """
    if not shape.has_text_frame:
        return False

    modified = False
    DATA_SLIDES = {1, 4, 6}
    if slide_count >= 12:
        DATA_SLIDES.add(10)
    is_data_slide = slide_idx is not None and slide_idx in DATA_SLIDES

    # Date range — replace within each run to preserve formatting (all shapes)
    if date_str:
        date_pattern = r'[A-Z][a-z]{2}\s+\d{1,2},\s+\d{4}\s*-\s*[A-Z][a-z]{2}\s+\d{1,2},\s+\d{4}'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                new = re.sub(date_pattern, date_str, run.text)
                if new != run.text:
                    run.text = new
                    modified = True

    col = get_column(shape.left)
    if not col:
        return modified  # only date replacement applies outside column zones

    # For each paragraph/run: if a heading key is found, replace the value;
    # otherwise clear stale template text (only on data-table slides, in data area).
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = run.text
            if not text.strip():
                continue
            heading_key = find_heading_key(text)
            if heading_key:
                report_val = get_value_from_report(heading_key, col, sections)
                formatted = format_val(report_val) if report_val is not None else '—'
                # Special handling for Previous Week column with "Actual: X" or "Target: Y" format
                if col == 'prev':
                    if 'Actual:' in text:
                        # Remove the word "Actual:" — keep just the value
                        prev_val = get_value_from_report(heading_key, 'prev', sections)
                        prev_fmt = format_val(prev_val) if prev_val is not None else '—'
                        new = re.sub(r'Actual:\s*', '', text)
                        new = re.sub(r'\b\d[\d,]*', prev_fmt, new, count=1)
                        if new != text:
                            run.text = new
                            modified = True
                        continue
                    elif text.strip().startswith('Target:'):
                        # Replace Target value with form's targetValue
                        target_val = get_value_from_report(heading_key, 'target', sections)
                        target_fmt = format_val(target_val) if target_val is not None else '—'
                        new = re.sub(r'Target:\s*[\d,+\-./day]+[\w/\s]*', f'Target: {target_fmt}', text)
                        if new != text:
                            run.text = new
                            modified = True
                        continue
                if ':' in text:
                    label_part = text.split(':')[0]
                    escaped_label = re.escape(label_part)
                    pattern = rf'({escaped_label}:\s*)([^;,\n]*?(?:\d[\d,+\-.]*)?)'
                    new = re.sub(pattern, lambda m: f'{m.group(1)}{formatted}', text, count=1, flags=re.IGNORECASE)
                    if new != text:
                        run.text = new
                        modified = True
                else:
                    new = re.sub(r'\b\d{1,3}(?:,\d{3})*\b', formatted, text, count=1)
                    if new != text:
                        run.text = new
                        modified = True
            elif is_data_slide and 1300000 <= shape.top <= 5800000:
                # No heading key — stale template text in data area, not in form
                run.text = ''
                modified = True

    return modified


def replace_narrative_for_dept(prs, sections, dept_name, slide_idx):
    """Replace narrative text on a specific slide with bullet-point notes from a department."""
    section = get_section_by_dept(sections, dept_name)
    notes = section.get('notes', '').strip() if section else ''
    # Clean up double-dots in sentences (data entry artifact)
    notes = re.sub(r'\.{2,}', '.', notes)
    bullets = '\n'.join(f"• {line.strip()}" for line in notes.split('\n') if line.strip()) if notes else ''
    slide = prs.slides[slide_idx]
    # Find the text box with the most content (the main narrative area)
    target = None
    target_len = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # Only consider sizable boxes in the content zone (below header, above footer)
        if shape.top < 700000 or shape.top > 6000000:
            continue
        if shape.width < 2000000 or shape.height < 500000:
            continue
        # Prefer boxes with existing content (the placeholder text)
        priority = len(text) if text else 0
        if priority > target_len:
            target_len = priority
            target = shape
    if target:
        # Preserve font properties from the first run before clearing
        kept_font = None
        first_para = target.text_frame.paragraphs[0]
        if first_para.runs:
            r0 = first_para.runs[0]
            kept_font = {
                'name': r0.font.name,
                'size': r0.font.size,
                'bold': r0.font.bold,
                'italic': r0.font.italic,
                'color': r0.font.color.rgb if r0.font.color and r0.font.color.type else None,
            }
        for para in target.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
        target.text_frame.paragraphs[0].runs[0].text = bullets
        # Restore font properties on the first run
        if kept_font:
            r0 = target.text_frame.paragraphs[0].runs[0]
            if kept_font['name']:
                r0.font.name = kept_font['name']
            if kept_font['size']:
                r0.font.size = kept_font['size']
            if kept_font['bold'] is not None:
                r0.font.bold = kept_font['bold']
            if kept_font['italic'] is not None:
                r0.font.italic = kept_font['italic']
            if kept_font['color']:
                from pptx.dml.color import RGBColor
                r0.font.color.rgb = kept_font['color']
        return True
    return False


SLIDE9_ROWS = [
    {'y_min': 3000000, 'y_max': 4000000, 'x_min': 0, 'x_max': 2500000, 'key': 'totalVendors'},
    {'y_min': 3000000, 'y_max': 4000000, 'x_min': 2500000, 'x_max': 6000000, 'key': 'totalMarketplaceProducts'},
    {'y_min': 3000000, 'y_max': 4000000, 'x_min': 6000000, 'x_max': 13000000, 'key': 'dailyAverageListings'},
]


def replace_summary_numbers(slide, sections, summary):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        col = get_column(shape.left)
        if not col:
            continue
        if re.match(r'^[\d,+\-./day\s]+$', text):
            for row in SLIDE9_ROWS:
                in_y = row['y_min'] <= shape.top <= row['y_max']
                in_x = row.get('x_min', 0) <= shape.left <= row.get('x_max', 99999999)
                if in_y and in_x:
                    heading_key = row['key']
                    # Closure KPI boxes always show current values
                    curr_val = get_value_from_report(heading_key, 'current', sections)
                    formatted = format_val(curr_val) if curr_val is not None else '—'
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            new = re.sub(r'[\d,+\-./day\s]+', formatted, run.text, count=1)
                            if new != run.text:
                                run.text = new
                    break


def build_12slide_template(template_path_9, output_path_12):
    """Build a 12-slide template from the 9-slide source by cloning and reordering slides.
    
    New structure (12 slides):
      0: Cover             (orig 0)
      1: BD Table          (orig 1)
      2: BD Charts         (orig 2)
      3: BD Narrative      (orig 3)
      4: Listing Table     (orig 4)
      5: Listing Charts    (orig 5)
      6: Listing Narrative (NEW — clone of orig 3, change header)
      7: QC Table          (orig 6)
      8: QC Charts         (NEW — clone of orig 2)
      9: QC Narrative      (NEW — clone of orig 3, change header)
     10: KPI Overview      (orig 7)
     11: Closure           (orig 8)
    """
    import zipfile
    from lxml import etree
    from copy import deepcopy
    import re as _re

    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    NS_CT = 'http://schemas.openxmlformats.org/package/2006/content-types'
    qn = lambda tag: '{%s}%s' % (NS_P, tag) if ':' not in tag else tag

    with zipfile.ZipFile(template_path_9, 'r') as z:
        all_data = {name: z.read(name) for name in z.namelist()}

    # ── Parse presentation.xml ──
    pres_xml = etree.fromstring(all_data['ppt/presentation.xml'])
    nsmap = pres_xml.nsmap
    sldIdLst = pres_xml.find(f'{{{NS_P}}}sldIdLst')
    orig_sld_ids = list(sldIdLst)
    orig_rIds = [s.get(f'{{{NS_R}}}id') for s in orig_sld_ids]  # e.g. ['rId2','rId3',...]

    # ── Parse presentation.xml.rels ──
    pres_rels_xml = etree.fromstring(all_data['ppt/_rels/presentation.xml.rels'])
    # findlast rId number
    max_rId_num = 0
    for child in pres_rels_xml:
        rid = child.get('Id', '')
        m = _re.search(r'rId(\d+)', rid)
        if m:
            max_rId_num = max(max_rId_num, int(m.group(1)))
    next_rId = max_rId_num + 1

    # Map orig rId -> target filename (e.g. 'rId2' -> 'slides/slide1.xml')
    orig_rId_to_target = {}
    for child in pres_rels_xml:
        rid = child.get('Id')
        target = child.get('Target', '')
        orig_rId_to_target[rid] = target

    # Read Content_Types to get existing slide overrides
    ct_xml = etree.fromstring(all_data.get('[Content_Types].xml', b''))
    next_slide_num = 10  # we'll add slides 10, 11, 12

    def _copy_slide_part(orig_rId):
        """Deep-copy a slide's XML and its rels file; return (slide_xml, rels_xml)."""
        target_rel = orig_rId_to_target[orig_rId]  # e.g. 'slides/slide3.xml'
        slide_path = f'ppt/{target_rel}'
        slide_xml = deepcopy(etree.fromstring(all_data[slide_path]))
        # Copy rels if present
        slide_name = target_rel.replace('slides/', '').replace('.xml', '')
        rels_path = f'ppt/slides/_rels/{slide_name}.xml.rels'
        rels_xml = deepcopy(etree.fromstring(all_data.get(rels_path, b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'))) if rels_path in all_data else None
        return slide_xml, rels_xml

    def _add_cloned_slide(src_rId, clone_xml, clone_rels_xml):
        """Register cloned slide parts into all_data and return its new rId."""
        nonlocal next_slide_num, all_data, pres_rels_xml, ct_xml, next_rId
        num = next_slide_num
        next_slide_num += 1
        new_file = f'slides/slide{num}.xml'
        new_rels = f'slides/_rels/slide{num}.xml.rels'

        # Add the slide XML
        all_data[f'ppt/{new_file}'] = etree.tostring(clone_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

        # Add the rels if present
        if clone_rels_xml is not None:
            all_data[f'ppt/{new_rels}'] = etree.tostring(clone_rels_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

        # Add content type
        override = etree.SubElement(ct_xml, f'{{{NS_CT}}}Override')
        override.set('PartName', f'/ppt/{new_file}')
        override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')

        # Add relationship in presentation.xml.rels
        rel_child = etree.SubElement(pres_rels_xml, 'Relationship')
        new_rId = f'rId{next_rId}'
        next_rId += 1
        rel_child.set('Id', new_rId)
        rel_child.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        rel_child.set('Target', new_file)

        # Ensure slideNum entry in [Content_Types].xml from template
        return new_rId

    def _get_slide_num_from_rId(rId):
        """Get slide number (1-based) from rId using presentation.xml.rels."""
        target = orig_rId_to_target.get(rId, '')
        m = _re.search(r'slide(\d+)', target)
        return int(m.group(1)) if m else None

    def _update_clone_header(slide_xml, old_dept, new_dept):
        """Replace department name in all text elements of a cloned slide's XML (case-insensitive)."""
        import re
        pattern = re.compile(re.escape(old_dept), re.IGNORECASE)
        for elem in slide_xml.iter():
            if elem.text:
                elem.text = pattern.sub(new_dept, elem.text)
            if elem.tail:
                elem.tail = pattern.sub(new_dept, elem.tail)

    # ── Clone the slides we need ──
    # orig 3 (BD Narrative) -> Listing Narrative clone
    rId_orig3 = orig_rIds[3]
    clone3_slide, clone3_rels = _copy_slide_part(rId_orig3)
    # Update header: "Business Development" -> "Listing"
    _update_clone_header(clone3_slide, 'Business Development', 'Listing')
    rId_clone3a = _add_cloned_slide(rId_orig3, clone3_slide, clone3_rels)

    # orig 2 (BD Charts) -> QC Charts clone
    rId_orig2 = orig_rIds[2]
    clone2_slide, clone2_rels = _copy_slide_part(rId_orig2)
    # Update header: "Business Development" -> "Quality Control" (any case)
    _update_clone_header(clone2_slide, 'Business Development', 'Quality Control')
    # Update chart labels for QC metrics
    _update_clone_header(clone2_slide, 'Total Vendors', 'Products Approved')
    _update_clone_header(clone2_slide, 'Total Verified Vendors', 'Products Rejected')
    rId_clone2 = _add_cloned_slide(rId_orig2, clone2_slide, clone2_rels)

    # orig 3 (BD Narrative) -> QC Narrative clone (another deep copy)
    clone3b_slide, clone3b_rels = _copy_slide_part(rId_orig3)
    # Update header: "Business Development" -> "Quality Control"
    _update_clone_header(clone3b_slide, 'Business Development', 'Quality Control')
    rId_clone3b = _add_cloned_slide(rId_orig3, clone3b_slide, clone3b_rels)

    # ── Reorder sldIdLst ──
    new_order = [
        ('orig', 0),  # Cover
        ('orig', 1),  # BD Table
        ('orig', 2),  # BD Charts
        ('orig', 3),  # BD Narrative
        ('orig', 4),  # Listing Table
        ('orig', 5),  # Listing Charts
        ('clone', rId_clone3a),  # Listing Narrative
        ('orig', 6),  # QC Table
        ('clone', rId_clone2),   # QC Charts
        ('clone', rId_clone3b),  # QC Narrative
        ('orig', 7),  # KPI Overview
        ('orig', 8),  # Closure
    ]
    # Clear existing sldId entries
    for s in orig_sld_ids:
        sldIdLst.remove(s)
    # Add in new order
    def _make_sldId(rId, idx):
        sldId = etree.SubElement(sldIdLst, f'{{{NS_P}}}sldId')
        sldId.set(f'{{{NS_R}}}id', rId)
        sldId.set('id', str(256 + idx))
        return sldId
    for idx, (typ, val) in enumerate(new_order):
        if typ == 'orig':
            rId = orig_rIds[val]
        else:
            rId = val
        _make_sldId(rId, idx)

    # ── Update presentation.xml in all_data ──
    all_data['ppt/presentation.xml'] = etree.tostring(pres_xml, xml_declaration=True, encoding='UTF-8', standalone=True)
    # Update presentation.xml.rels
    all_data['ppt/_rels/presentation.xml.rels'] = etree.tostring(pres_rels_xml, xml_declaration=True, encoding='UTF-8', standalone=True)
    # Update [Content_Types].xml
    all_data['[Content_Types].xml'] = etree.tostring(ct_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

    # ── Write output PPTX ──
    import shutil
    output_path_12 = str(output_path_12)
    os.makedirs(os.path.dirname(output_path_12) or '.', exist_ok=True)
    with zipfile.ZipFile(output_path_12, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in sorted(all_data.items()):
            zout.writestr(name, data)

    print(f"✓ 12-slide template saved to {output_path_12}")
    return output_path_12


def generate_pptx(data):
    from pptx import Presentation
    from pptx.oxml.ns import qn

    template_path = data['template_path']
    report = data.get('report', {})
    sections = report.get('sections', [])
    summary = report.get('summary', {})
    date_str = report.get('nepaliDate', '')

    # ── Build 12-slide template if using the 9-slide source ──
    use_12slide = data.get('use_12slide', True)
    if use_12slide:
        cache_path = '/tmp/opencode/template_12slide.pptx'
        if not os.path.exists(cache_path):
            os.makedirs('/tmp/opencode', exist_ok=True)
            build_12slide_template(template_path, cache_path)
        template_path = cache_path

    prs = Presentation(template_path)

    slide_count = len(prs.slides)

    # ── Replace text in all shapes across all slides ──
    for slide_idx, slide in enumerate(prs.slides):
        # Build heading_key map for Previous Week "Target:" shapes (by position)
        prev_actual_positions = {}  # (left, top) -> heading_key for "Actual:" shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if get_column(shape.left) != 'prev':
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if 'Actual:' in text:
                hk = find_heading_key(text)
                if hk:
                    prev_actual_positions[(shape.left, shape.top)] = hk
            elif text.startswith('Target:'):
                # Clear Target: shapes in the Previous Week column — only actual data here
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
        
        # Process all other shapes normally
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            # Skip shapes already handled above (Target: in prev column)
            if get_column(shape.left) == 'prev' and text.startswith('Target:'):
                continue
            if text:
                replace_label_values_in_shape(shape, sections, summary, date_str, slide_idx, slide_count)

        # Narrative slides (indices differ based on template version)
        if slide_idx == 3:
            replace_narrative_for_dept(prs, sections, 'Business Development', 3)
        if slide_count >= 12:
            if slide_idx == 6:
                replace_narrative_for_dept(prs, sections, 'Listing', 6)
            if slide_idx == 9:
                replace_narrative_for_dept(prs, sections, 'Quality Control', 9)
            if slide_idx == 11:
                replace_summary_numbers(slide, sections, summary)
        else:
            # Old 9-slide template: summary on slide 8
            if slide_idx == 8:
                replace_summary_numbers(slide, sections, summary)

    # ── Chart images ──────────────────────────────────
    BAR_LABELS = ['Previous\nWeek', 'Current\nWeek', 'Next Week\nTarget']

    def slide(n):
        """Safely get slide by index, returns None if out of range."""
        return prs.slides[n] if n < len(prs.slides) else None

    def shape_inches(shape):
        """Return (width_in, height_in) from a shape's EMU dimensions."""
        return (shape.width / 914400, shape.height / 914400)

    # Slide 2 (index 2): BD Progress Charts (same in both 9 and 12 slide mode)
    slide_bd_charts = slide(2)
    if slide_bd_charts:
        left_chart1 = find_picture_by_position(slide_bd_charts, 411480, 1824770)
        right_chart1 = find_picture_by_position(slide_bd_charts, 6420700, 1866850)

        if left_chart1:
            w, h = shape_inches(left_chart1)
            chart_vals = get_chart_data_3bars(sections, 'Business Development', 'totalVendors')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'Total Vendors', fig_width=w, fig_height=h)
            _replace_image_with_new_part(left_chart1, buf.read())

        if right_chart1:
            w, h = shape_inches(right_chart1)
            chart_vals = get_chart_data_3bars(sections, 'Business Development', 'verifiedVendors')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'Total Verified Vendors', fig_width=w, fig_height=h)
            _replace_image_with_new_part(right_chart1, buf.read())

    # Slide 5 (index 5): Listing Progress Charts (same in both modes)
    slide_listing_charts = slide(5)
    if slide_listing_charts:
        LEFT_CHART_POS = (713232, 1600200)  # fallback for 9-slide
        RIGHT_CHART_POS = (6812280, 1600200)
        left_chart2 = find_picture_by_position(slide_listing_charts, *LEFT_CHART_POS)
        if left_chart2:
            w, h = shape_inches(left_chart2)
            chart_vals = get_chart_data_3bars(sections, 'Listing', 'totalMarketplaceProducts')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'Marketplace Products', fig_width=w, fig_height=h)
            _replace_image_with_new_part(left_chart2, buf.read())

        right_chart2a = find_picture_by_position(slide_listing_charts, *RIGHT_CHART_POS)
        if right_chart2a:
            w, h = shape_inches(right_chart2a)
            chart_vals = get_chart_data_3bars(sections, 'Listing', 'totalSpecificationsAdded')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'Specifications Added', fig_width=w, fig_height=h)
            _replace_image_with_new_part(right_chart2a, buf.read())

    # Slide 8 (12-slide) or slide 8 (9-slide): QC Charts
    slide_qc_charts = slide(8)
    if slide_qc_charts:
        qc_metrics = [
            ('productsApproved', 'QC – Products Approved'),
            ('productsRejected', 'QC – Products Rejected'),
            ('productsPending', 'QC – Products Pending'),
        ]
        left_positions = [(411480, 1824770), (6420700, 1866850)]
        old_qc_positions = [
            (457200, 1508760),
            (457200, 1508760),
            (457200, 1508760),
        ]
        for i, (metric, title) in enumerate(qc_metrics):
            placeholder = None
            if i < len(left_positions):
                placeholder = find_picture_by_position(slide_qc_charts, *left_positions[i])
            if not placeholder and i < len(old_qc_positions):
                placeholder = find_picture_by_position(slide_qc_charts, *old_qc_positions[i])
            if placeholder:
                w, h = shape_inches(placeholder)
                chart_vals = get_chart_data_3bars(sections, 'Quality Control', metric)
                buf = create_bar_chart(chart_vals, BAR_LABELS, title, fig_width=w, fig_height=h)
                _replace_image_with_new_part(placeholder, buf.read())

    # Slide 10 (index 10): KPI Overview (old index 7 in 9-slide, or old index 8)
    slide_kpi = slide(10)
    if slide_kpi:
        left_chart3 = find_picture_by_position(slide_kpi, 685800, 1508760)
        if left_chart3:
            w, h = shape_inches(left_chart3)
            chart_vals = get_chart_data_3bars(sections, 'Business Development', 'totalVendors')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'BD KPI — Total Vendors', fig_width=w, fig_height=h)
            _replace_image_with_new_part(left_chart3, buf.read())

        right_chart3 = find_picture_by_position(slide_kpi, 6446520, 1508760)
        if right_chart3:
            w, h = shape_inches(right_chart3)
            chart_vals = get_chart_data_3bars(sections, 'Listing', 'totalMarketplaceProducts')
            buf = create_bar_chart(chart_vals, BAR_LABELS, 'Listing — Marketplace Products', fig_width=w, fig_height=h)
            _replace_image_with_new_part(right_chart3, buf.read())

    # ── Save ──
    output_path = data.get('output_path', 'output.pptx')
    prs.save(str(output_path))
    print(f"✓ Saved to {output_path}")


# ── PDF Generation ─────────────────────────────────────

def _register_fonts(pdf):
    regular = FONT_DIR / 'DejaVuSans.ttf'
    bold = FONT_DIR / 'DejaVuSans-Bold.ttf'
    if regular.exists():
        pdf.add_font('DejaVu', '', str(regular), uni=True)
    if bold.exists():
        pdf.add_font('DejaVu', 'B', str(bold), uni=True)


def _pdf_header(pdf, text, size=16):
    pdf.set_font('DejaVu', 'B', size)
    pdf.set_text_color(180, 21, 49)
    pdf.cell(0, 10, text, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)


def _add_cover_page(pdf, report, date_str):
    pdf.add_page()
    _register_fonts(pdf)
    pdf.ln(35)
    pdf.set_font('DejaVu', 'B', 26)
    pdf.set_text_color(180, 21, 49)
    pdf.cell(0, 14, 'Nepalcan.com', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.set_font('DejaVu', 'B', 16)
    pdf.set_text_color(60, 60, 60)
    pdf.cell(0, 10, 'Weekly Review', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(8)
    pdf.set_font('DejaVu', '', 12)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 8, 'Department-Wise Performance Analysis', align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    pdf.set_font('DejaVu', '', 11)
    pdf.cell(0, 7, date_str, align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font('DejaVu', '', 10)
    pdf.set_text_color(130, 130, 130)
    pdf.cell(0, 7, 'Executive Brief', align='C', new_x="LMARGIN", new_y="NEXT")


def _add_dept_table(pdf, sections, dept_name):
    section = get_section_by_dept(sections, dept_name)
    if not section:
        return
    pdf.add_page()
    _register_fonts(pdf)
    _pdf_header(pdf, dept_name.upper())
    pdf.set_font('DejaVu', '', 7)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 5, 'PREVIOUS Week    |    CURRENT STATUS    |    NEXT WEEK TARGET', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)

    cols = [85, 30, 30, 30]
    headers = ['Metric', 'Previous', 'Current', 'Target']

    pdf.set_font('DejaVu', 'B', 8)
    pdf.set_fill_color(244, 244, 245)
    for i, h in enumerate(headers):
        align = 'L' if i == 0 else 'R'
        pdf.cell(cols[i], 7, h, border=1, fill=True, align=align)
    pdf.ln()

    pdf.set_font('DejaVu', '', 8)
    for v in section.get('values', []):
        name = v.get('headingName', '')
        pdf.cell(cols[0], 6, name, border=1, align='L')
        pdf.cell(cols[1], 6, format_val(v.get('previousValue')), border=1, align='R')
        pdf.set_fill_color(254, 242, 242)
        pdf.cell(cols[2], 6, format_val(v.get('currentValue')), border=1, align='R', fill=True)
        pdf.set_text_color(217, 119, 6)
        pdf.cell(cols[3], 6, format_val(v.get('targetValue')), border=1, align='R')
        pdf.set_text_color(60, 60, 60)
        pdf.ln()


def _add_chart_pdf_page(pdf, sections, dept_name, charts, extra_chart=None):
    pdf.add_page()
    _register_fonts(pdf)
    _pdf_header(pdf, f'{dept_name.upper()} - PROGRESS CHARTS', 14)
    pdf.ln(2)

    all_charts = list(charts)
    if extra_chart:
        all_charts.append(extra_chart)

    n = len(all_charts)
    cw = 85 if n <= 2 else 57
    gap = 4 if n <= 2 else 3
    total_w = n * cw + (n - 1) * gap
    x_start = (210 - total_w) / 2
    ch = 55 if n <= 2 else 40

    BAR_LABELS_SHORT = ['Prev', 'Current', 'Target']
    for i, (chart_dept, key, title) in enumerate(all_charts):
        x = x_start + i * (cw + gap)
        y = pdf.get_y()
        buf = create_bar_chart(
            get_chart_data_3bars(sections, chart_dept, key),
            BAR_LABELS_SHORT, title, dpi=150
        )
        pdf.image(buf, x=x, y=y, w=cw, h=ch)


def _add_narrative_pdf_page(pdf, sections, dept_name):
    section = get_section_by_dept(sections, dept_name)
    pdf.add_page()
    _register_fonts(pdf)
    _pdf_header(pdf, f'{dept_name} Report', 14)
    notes = section.get('notes', '').strip() if section else ''
    if notes:
        pdf.set_font('DejaVu', '', 10)
        pdf.set_text_color(50, 50, 50)
        for line in notes.split('\n'):
            line = line.strip()
            if line:
                pdf.set_x(20)
                pdf.cell(0, 7, f'  •  {line}', new_x="LMARGIN", new_y="NEXT")
    else:
        pdf.set_font('DejaVu', '', 10)
        pdf.set_text_color(180, 180, 180)
        pdf.cell(0, 7, 'No notes provided.', new_x="LMARGIN", new_y="NEXT")


def _add_summary_pdf_page(pdf, sections, summary):
    pdf.add_page()
    _register_fonts(pdf)
    _pdf_header(pdf, 'REPORT CLOSURE', 14)
    pdf.ln(3)
    pdf.set_font('DejaVu', 'B', 20)
    pdf.set_text_color(50, 50, 50)
    pdf.cell(0, 12, 'Thank You', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    date_str = ''
    if report := None:
        pass
    pdf.set_font('DejaVu', '', 9)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 6, 'Summary', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    kpis = [
        ('Total Vendors', 'totalVendors', 'Business Development'),
        ('Total Marketplace Products', 'totalMarketplaceProducts', 'Listing'),
        ('Daily Average Listings', 'dailyAverageListings', 'Listing'),
    ]
    cw = 56
    gap = 4
    total_w = 3 * cw + 2 * gap
    x_start = (210 - total_w) / 2

    for i, (label, key, dept) in enumerate(kpis):
        x = x_start + i * (cw + gap)
        y = pdf.get_y()
        val = summary.get(key) if summary else None
        if val is None:
            section = get_section_by_dept(sections, dept)
            val = get_value_from_section(section, key, 'current')
        pdf.set_draw_color(200, 200, 200)
        pdf.set_line_width(0.3)
        pdf.rect(x, y, cw, 30, style='D')
        pdf.set_xy(x + 2, y + 4)
        pdf.set_font('DejaVu', 'B', 16)
        pdf.set_text_color(30, 30, 30)
        pdf.cell(cw - 4, 10, format_val(val) if val is not None else '—', align='C', new_x="LMARGIN", new_y="NEXT")
        pdf.set_x(x + 2)
        pdf.set_font('DejaVu', '', 7)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(cw - 4, 5, label, align='C')


def generate_pdf(data):
    report = data.get('report', {})
    sections = report.get('sections', [])
    summary = report.get('summary', {})
    date_str = report.get('nepaliDate', '')

    pdf = FPDF('P', 'mm', 'A4')
    pdf.set_auto_page_break(auto=True, margin=18)

    _add_cover_page(pdf, report, date_str)
    _add_dept_table(pdf, sections, 'Business Development')
    _add_chart_pdf_page(pdf, sections, 'Business Development', [
        ('Business Development', 'totalVendors', 'Total Vendors'),
        ('Business Development', 'verifiedVendors', 'Total Verified Vendors'),
    ])
    _add_narrative_pdf_page(pdf, sections, 'Business Development')
    _add_dept_table(pdf, sections, 'Listing')
    _add_chart_pdf_page(pdf, sections, 'Listing', [
        ('Listing', 'totalMarketplaceProducts', 'Marketplace Products'),
        ('Listing', 'totalSpecificationsAdded', 'Specifications Added'),
    ])
    _add_narrative_pdf_page(pdf, sections, 'Listing')
    _add_dept_table(pdf, sections, 'Quality Control')
    _add_chart_pdf_page(pdf, sections, 'Quality Control', [
        ('Quality Control', 'productsApproved', 'Products Approved'),
        ('Quality Control', 'productsRejected', 'Products Rejected'),
        ('Quality Control', 'productsPending', 'Products Pending'),
    ])
    _add_narrative_pdf_page(pdf, sections, 'Quality Control')
    _add_chart_pdf_page(pdf, sections, 'Business Development', [
        ('Business Development', 'totalVendors', 'BD KPI - Total Vendors'),
    ], extra_chart=('Listing', 'totalMarketplaceProducts', 'Listing - Marketplace Products'))
    _add_summary_pdf_page(pdf, sections, summary)

    output_path = data.get('pdf_output_path', 'output.pdf')
    pdf.output(str(output_path))
    print(f"✓ Saved PDF to {output_path}")


def main():
    parser = argparse.ArgumentParser(description='Generate weekly report PPTX and/or PDF')
    parser.add_argument('--data', required=True, help='Path to JSON data file')
    parser.add_argument('--output', help='Output PPTX path')
    parser.add_argument('--pdf', help='Output PDF path')
    parser.add_argument('--debug-slide', type=int, help='Print debug info for slide N (0-indexed)')
    args = parser.parse_args()

    if not args.output and not args.pdf:
        parser.error('At least one of --output or --pdf must be specified')

    # Debug mode: print slide info and exit
    if args.debug_slide is not None:
        from pptx import Presentation
        data_path = args.data or 'NEPALCAN.COM-2026-WK-20 (1).pptx'
        with open(data_path) as f:
            data = json.load(f)
        template_path = data['template_path']
        prs = Presentation(template_path)
        slide = prs.slides[args.debug_slide]
        print(f"\n=== Slide {args.debug_slide} Debug Info ===")
        for i, shape in enumerate(slide.shapes):
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if hasattr(shape, 'image') and shape.shape_type == 13:
                print(f"[{i}] IMAGE: left={left}, top={top}, width={width}, height={height}")
            elif shape.has_text_frame:
                text_preview = shape.text_frame.text[:50].replace('\n', ' ') if len(shape.text_frame.text) > 50 else shape.text_frame.text.replace('\n', ' ')
                print(f"[{i}] TEXT: left={left}, top={top}, width={width}, height={height}, text='{text_preview}'")
            else:
                print(f"[{i}] OTHER: left={left}, top={top}, width={width}, height={height}, shape_type={shape.shape_type}")
        return

    with open(args.data) as f:
        data = json.load(f)

    if args.output:
        data['output_path'] = args.output
        generate_pptx(data)

    if args.pdf:
        data['pdf_output_path'] = args.pdf
        generate_pdf(data)


if __name__ == '__main__':
    if '--data' not in ' '.join(sys.argv) and '--debug-slide' not in ' '.join(sys.argv):
        print("Usage: python generate_pptx_report.py --data <json> --output <pptx> [--pdf <pdf>]")
        print("       python generate_pptx_report.py --data <json> --debug-slide 1")
        sys.exit(1)
    main()
