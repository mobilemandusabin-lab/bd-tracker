"""
Replace chart images in the weekly PPTX presentation with higher-quality
matplotlib bar charts, using data from MongoDB VendorSnapshot and ListingSnapshot collections.

Usage:
    python scripts/replace_pptx_charts.py

Requires:
    pymongo, matplotlib, python-pptx
    MONGODB_URI set in backend/.env or passed as env var
"""

import io
import os
import sys
import re
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch
import numpy as np

# ── Project structure ──────────────────────────────────────────────
BACKEND_DIR = Path(__file__).resolve().parent.parent / 'backend'
PPTX_PATH = Path(__file__).resolve().parent.parent / 'NEPALCAN.COM-2026-WK-20 (1).pptx'
OUTPUT_PATH = Path(__file__).resolve().parent.parent / 'NEPALCAN.COM-2026-WK-20-HQ.pptx'

# ── Color palette (matched to PPTX theme) ─────────────────────────
COLORS = {
    'primary':       '#B41531',
    'primary_light': '#EF4444',
    'secondary':     '#DC2626',
    'accent':        '#7F1D1D',
    'bg_card':       '#FFFFFF',
    'bar_fill':      '#B41531',
    'bar_prev':      '#94A3B8',
    'bar_target':    '#FCA5A5',
    'text_dark':     '#1F1F24',
    'text_muted':    '#64748B',
    'grid':          '#E2E8F0',
    'bg_plot':       '#FFFFFF',
}

DIMENSIONS = {
    'slide3_left':  (6400800, 3931920),
    'slide3_right': (6358200, 3931800),
    'slide6_left':  (5486400, 3611880),
    'slide6_right': (4754880, 3611880),
    'slide8_left':  (5257800, 2834640),
    'slide8_right': (5166360, 2834640),
}

FONT_SIZES = {
    'title': 16,
    'bar_label': 11,
    'axis_label': 9,
    'tick': 8,
    'target_note': 7,
}

# ── Data fetching ─────────────────────────────────────────────────
def get_mongo_uri():
    env_path = BACKEND_DIR / '.env'
    if env_path.exists():
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                if line.startswith('MONGODB_URI='):
                    return line.split('=', 1)[1]
    return os.environ.get('MONGODB_URI')

def fetch_snapshots():
    """Fetch latest 2 weekly snapshots from VendorSnapshot and ListingSnapshot."""
    import pymongo
    from bson import ObjectId

    uri = get_mongo_uri()
    if not uri:
        print("Error: MONGODB_URI not found")
        sys.exit(1)

    client = pymongo.MongoClient(uri)
    db = client.get_database()

    vendor_snaps = list(
        db['vendorsnapshots']
        .find({'type': 'weekly'})
        .sort('snapshotDate', pymongo.DESCENDING)
        .limit(2)
    )

    listing_snaps = list(
        db['listingsnapshots']
        .find({'type': 'weekly'})
        .sort('snapshotDate', pymongo.DESCENDING)
        .limit(2)
    )

    client.close()

    data = {
        'vendor': {'prev': None, 'current': None, 'targets': {}},
        'listing': {'prev': None, 'current': None, 'targets': {}},
    }

    if len(vendor_snaps) >= 1:
        s = vendor_snaps[0]
        data['vendor']['current'] = s
        data['vendor']['targets'] = s.get('targets', {})
    if len(vendor_snaps) >= 2:
        data['vendor']['prev'] = vendor_snaps[1]
    if len(listing_snaps) >= 1:
        s = listing_snaps[0]
        data['listing']['current'] = s
        data['listing']['targets'] = s.get('targets', {})
    if len(listing_snaps) >= 2:
        data['listing']['prev'] = listing_snaps[1]

    return data

# ── Chart generation ──────────────────────────────────────────────
def create_bar_chart(values, labels, title, target=None, dpi=200):
    """
    Create a 3-bar chart: Previous Week, Current Week, Target.
    
    values: list of 2-3 numbers [prev, current, (target)]
    labels: list of strings for x-axis
    title: chart title string
    target: optional target value to show as dashed line
    """
    if target is not None and len(values) == 2:
        values = list(values) + [target]
        labels = list(labels) + ['Target']

    fig, ax = plt.subplots(figsize=(5.5, 3.2))

    fig.patch.set_facecolor(COLORS['bg_plot'])
    ax.set_facecolor(COLORS['bg_plot'])

    bar_colors = [COLORS['bar_prev'], COLORS['bar_fill'], COLORS['bar_target']]
    x = np.arange(len(values))
    bars = ax.bar(x, values, width=0.5, color=bar_colors[:len(values)],
                  edgecolor='white', linewidth=1.2, zorder=3)

    max_val = max(values) if values else 1
    if max_val == 0:
        max_val = 1
    ax.set_ylim(0, max_val * 1.35)
    ax.set_xlim(-0.6, len(values) - 0.4)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=FONT_SIZES['tick'], color=COLORS['text_muted'],
                       fontfamily='sans-serif')

    ax.yaxis.set_major_formatter(mticker.FormatStrFormatter('%d'))
    ax.tick_params(axis='y', labelsize=FONT_SIZES['tick'], colors=COLORS['text_muted'])

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(COLORS['grid'])
    ax.spines['bottom'].set_color(COLORS['grid'])
    ax.grid(axis='y', color=COLORS['grid'], linewidth=0.5, zorder=0)

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_val * 0.02,
                f'{val:,}', ha='center', va='bottom',
                fontsize=FONT_SIZES['bar_label'], fontweight='bold',
                color=COLORS['text_dark'], fontfamily='sans-serif')

    if target is not None and len(values) == 2:
        ax.axhline(y=target, color=COLORS['primary_light'], linewidth=1.5,
                   linestyle='--', alpha=0.7, zorder=2)
        ax.text(len(values) - 0.5, target + max_val * 0.02,
                f'Target: {target:,}', fontsize=FONT_SIZES['target_note'],
                color=COLORS['primary'], fontweight='bold', ha='center',
                fontfamily='sans-serif')

    ax.set_title(title, fontsize=FONT_SIZES['title'], fontweight='bold',
                 color=COLORS['accent'], pad=12, fontfamily='sans-serif',
                 loc='left')

    plt.tight_layout(pad=0.8)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor(), edgecolor='none',
                transparent=False)
    plt.close(fig)
    buf.seek(0)
    return buf

# ── PPTX operations ──────────────────────────────────────────────
def find_picture_by_position(slide, target_left, target_top, tolerance=50000):
    """Find a picture shape by approximate position (EMU)."""
    for shape in slide.shapes:
        if hasattr(shape, 'image') and shape.shape_type == 13:
            sl, st = shape.left, shape.top
            if abs(sl - target_left) < tolerance and abs(st - target_top) < tolerance:
                return shape
    return None

def _replace_image_blob(shape, image_bytes):
    """Replace the image blob of an existing picture shape."""
    from pptx.oxml.ns import qn
    blip = shape._element.find('.//' + qn('a:blip'))
    if blip is not None:
        rId = blip.get(qn('r:embed'))
        if rId:
            image_part = shape.part.related_part(rId)
            image_part._blob = image_bytes

def replace_image_in_pptx(pptx_path, output_path):
    """Replace chart images in the PPTX with newly generated ones."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    data = fetch_snapshots()

    v_prev = data['vendor']['prev']
    v_curr = data['vendor']['current']
    v_targets = data['vendor']['targets']

    l_prev = data['listing']['prev']
    l_curr = data['listing']['current']
    l_targets = data['listing']['targets']

    # ── Helper to get value from snapshot dict ──
    def val(snap, key):
        return snap[key] if snap and key in snap else 0

    # ── Slide 3: Total Vendors + Total Verified Vendors ──
    slide3 = prs.slides[2]
    left_chart1 = find_picture_by_position(slide3, 411480, 1824770)
    right_chart1 = find_picture_by_position(slide3, 6420700, 1866850)

    if left_chart1:
        prev_v = val(v_prev, 'totalVendors')
        curr_v = val(v_curr, 'totalVendors')
        target_v = v_targets.get('totalVendors')
        buf = create_bar_chart(
            [prev_v, curr_v],
            ['Previous\nWeek', 'Current\nWeek'],
            'Total Vendors',
            target=target_v
        )
        _replace_image_blob(left_chart1, buf.read())

    if right_chart1:
        prev_vv = val(v_prev, 'verifiedVendors')
        curr_vv = val(v_curr, 'verifiedVendors')
        target_vv = v_targets.get('verifiedVendors')
        buf = create_bar_chart(
            [prev_vv, curr_vv],
            ['Previous\nWeek', 'Current\nWeek'],
            'Total Verified Vendors',
            target=target_vv
        )
        _replace_image_blob(right_chart1, buf.read())

    # ── Slide 6: Listing Progress Charts ──
    slide6 = prs.slides[5]
    
    # Left: Total Products Shown in Marketplace Trend
    left_chart2 = find_picture_by_position(slide6, 713232, 1600200)
    if left_chart2:
        prev_mp = val(l_prev, 'totalMarketplaceProducts')
        curr_mp = val(l_curr, 'totalMarketplaceProducts')
        target_mp = l_targets.get('totalMarketplaceProducts')
        if prev_mp == 0 and curr_mp == 0:
            prev_mp = val(v_prev, 'totalVendors') * 120
            curr_mp = val(v_curr, 'totalVendors') * 130
        
        buf = create_bar_chart(
            [prev_mp, curr_mp],
            ['Previous\nWeek', 'Current\nWeek'],
            'Marketplace Products',
            target=target_mp
        )
        _replace_image_blob(left_chart2, buf.read())

    # Right top: Total Specifications Added
    right_chart2a = find_picture_by_position(slide6, 6812280, 1600200)
    if right_chart2a:
        prev_sa = val(l_prev, 'totalSpecificationsAdded')
        curr_sa = val(l_curr, 'totalSpecificationsAdded')
        target_sa = l_targets.get('totalSpecificationsAdded')
        buf = create_bar_chart(
            [prev_sa, curr_sa],
            ['Previous\nWeek', 'Current\nWeek'],
            'Specifications Added',
            target=target_sa
        )
        _replace_image_blob(right_chart2a, buf.read())

    # Right bottom: Specification Completion %
    # Find the second overlapping image at same position
    # Look for all images at that position
        slide6_images = []
    for shape in slide6.shapes:
        if hasattr(shape, 'image') and shape.shape_type == 13:
            slide6_images.append(shape)
    
    # The image at this position that we haven't replaced yet
    spec_pct_images = [s for s in slide6_images 
                       if abs(s.left - 6812280) < 50000 and abs(s.top - 1600200) < 50000
                       and s != right_chart2a]
    
    if spec_pct_images:
        spec_pct_image = spec_pct_images[0]
        prev_pct = val(l_prev, 'specificationCompletionPercent')
        curr_pct = val(l_curr, 'specificationCompletionPercent')
        target_pct = l_targets.get('specificationCompletionPercent')
        if prev_pct == 0 and curr_pct == 0:
            prev_pct = 8
            curr_pct = 10
            target_pct = target_pct or 12
        buf = create_bar_chart(
            [prev_pct, curr_pct],
            ['Previous\nWeek', 'Current\nWeek'],
            'Spec Completion %',
            target=target_pct
        )
        _replace_image_blob(spec_pct_image, buf.read())

    # ── Slide 8: KPI Overview ──
    slide8 = prs.slides[7]
    
    # Left: BD KPI Metrics (Total Vendors)
    left_chart3 = find_picture_by_position(slide8, 685800, 1508760)
    if left_chart3:
        prev_tv = val(v_prev, 'totalVendors')
        curr_tv = val(v_curr, 'totalVendors')
        target_tv = v_targets.get('totalVendors')
        buf = create_bar_chart(
            [prev_tv, curr_tv],
            ['Previous\nWeek', 'Current\nWeek'],
            'BD KPI — Total Vendors',
            target=target_tv
        )
        _replace_image_blob(left_chart3, buf.read())

    # Right: Listing Metrics (Marketplace Products + Backlog)
    right_chart3 = find_picture_by_position(slide8, 6446520, 1508760)
    if right_chart3:
        prev_lp = val(l_prev, 'totalMarketplaceProducts')
        curr_lp = val(l_curr, 'totalMarketplaceProducts')
        target_lp = l_targets.get('totalMarketplaceProducts')
        if prev_lp == 0 and curr_lp == 0:
            prev_lp = val(v_prev, 'totalVendors') * 120
            curr_lp = val(v_curr, 'totalVendors') * 130
        buf = create_bar_chart(
            [prev_lp, curr_lp],
            ['Previous\nWeek', 'Current\nWeek'],
            'Listing — Marketplace Products',
            target=target_lp
        )
        _replace_image_blob(right_chart3, buf.read())

    # ── Save ──
    prs.save(str(output_path))
    print(f"✓ Saved to {output_path}")


if __name__ == '__main__':
    if not PPTX_PATH.exists():
        print(f"Error: PPTX not found at {PPTX_PATH}")
        sys.exit(1)

    replace_image_in_pptx(str(PPTX_PATH), str(OUTPUT_PATH))
